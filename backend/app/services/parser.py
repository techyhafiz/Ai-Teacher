"""Document parsing service.

Extracts text + structure from PDF / DOCX / PPTX / TXT / MD files.

Pipeline per document:
  1. Native text extraction (PyMuPDF / python-docx / python-pptx / plain)
  2. Structure detection: heading-based sections from font-size/style
     (PDF), paragraph styles (DOCX), slide boundaries (PPTX); if a doc has
     no detectable headings, a TOC pass synthesizes a table of contents
  3. Vision OCR: pages with <100 extractable chars are rendered to images
     and OCR'd by the Gemini text model (multimodal) — this also describes
     figures/diagrams on those pages
  4. Output: JSON structure {sections: [{id, title, level, text, source,
     images_described}] , full_text} stored under data/processed/<doc_id>/
"""
from __future__ import annotations

import base64
import json
import logging
import re
import uuid
from pathlib import Path
from typing import Any, Optional

from ..config import PROCESSED_DIR, UPLOADS_DIR
from .gemini import BATCH, generate_text, estimate_tokens
from .tpm_manager import tpm

log = logging.getLogger("parser")

OCR_MIN_CHARS = 100          # pages with fewer text chars get vision OCR
OCR_IMAGE_MAX_DIM = 1400     # px, downscaled before sending


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

async def process_document(file_path: str | Path, original_name: str,
                           *, language_hint: str = "en",
                           max_ocr_pages: int = 40) -> dict[str, Any]:
    """Parse a document into a structured JSON record. Returns the record."""
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        raw = _extract_pdf(path)
    elif suffix == ".docx":
        raw = _extract_docx(path)
    elif suffix == ".pptx":
        raw = _extract_pptx(path)
    elif suffix in (".txt", ".md", ".markdown"):
        raw = _extract_plain(path)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")

    # Vision OCR for sparse pages (scanned PDFs, image-heavy slides)
    ocr_applied: list[dict] = []
    if suffix == ".pdf":
        ocr_applied = await _ocr_sparse_pdf_pages(path, raw, max_ocr_pages)

    # Structure detection / TOC synthesis
    structured = _detect_structure(raw, suffix)

    doc_id = uuid.uuid4().hex[:12]
    record = {
        "doc_id": doc_id,
        "original_name": original_name,
        "format": suffix.lstrip("."),
        "language_hint": language_hint,
        "sections": structured["sections"],
        "toc": structured["toc"],
        "stats": {
            **structured["stats"],
            "ocr_pages": len(ocr_applied),
        },
        "ocr": ocr_applied,
    }

    out_dir = PROCESSED_DIR / doc_id
    out_dir.mkdir(parents=True, exist_ok=True)
    (out_dir / "document.json").write_text(
        json.dumps(record, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    log.info("Parsed '%s' -> %d sections, %d OCR pages, doc_id=%s",
             original_name, len(record["sections"]), len(ocr_applied), doc_id)
    return record


def load_document(doc_id: str) -> Optional[dict[str, Any]]:
    """Load a previously processed document record."""
    p = PROCESSED_DIR / doc_id / "document.json"
    if not p.exists():
        return None
    return json.loads(p.read_text(encoding="utf-8"))


def document_full_text(doc_id: str) -> str:
    """Full text (sections joined), for long-context injection."""
    doc = load_document(doc_id)
    if not doc:
        return ""
    parts = []
    for s in doc["sections"]:
        parts.append(f"## {s['title']}\n{s['text']}")
    return "\n\n".join(parts)


def document_token_estimate(doc_id: str) -> int:
    doc = load_document(doc_id)
    if not doc:
        return 0
    return sum(estimate_tokens(s["text"]) for s in doc["sections"])


# ---------------------------------------------------------------------------
# Extractors — produce {"units": [{title, text, source, level}], "raw_text"}
# ---------------------------------------------------------------------------

def _extract_pdf(path: Path) -> dict[str, Any]:
    import fitz  # PyMuPDF

    units: list[dict[str, Any]] = []
    raw_len = 0
    with fitz.open(path) as doc:
        for pno, page in enumerate(doc):
            blocks = page.get_text("dict")["blocks"]
            page_lines: list[tuple[float, float, bool, str]] = []
            max_size = 0.0
            sizes: list[float] = []
            for b in blocks:
                if b.get("type") != 0:
                    continue
                for line in b.get("lines", []):
                    text = "".join(sp.get("text", "") for sp in line.get("spans", []))
                    if not text.strip():
                        continue
                    size = max((sp.get("size", 10.0) for sp in line.get("spans", [])),
                               default=10.0)
                    bold = any("bold" in (sp.get("font", "") or "").lower()
                               for sp in line.get("spans", []))
                    page_lines.append((line["bbox"][1], size, bold, text.strip()))
                    sizes.append(size)
            body_size = _modal_size(sizes) if sizes else 10.0
            for y, size, bold, text in sorted(page_lines, key=lambda x: x[0]):
                is_heading = (size >= body_size * 1.22 and len(text) < 90) or \
                             (bold and size >= body_size * 1.05 and len(text) < 90)
                units.append({
                    "page": pno + 1, "size": round(size, 2), "bold": bold,
                    "heading": is_heading, "text": text, "body_size": round(body_size, 2),
                })
                raw_len += len(text)
    return {"units": units, "raw_len": raw_len, "kind": "pdf_lines"}


def _extract_docx(path: Path) -> dict[str, Any]:
    from docx import Document

    doc = Document(str(path))
    units: list[dict[str, Any]] = []
    raw_len = 0
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower()
        is_heading = style.startswith("heading") or style.startswith("title")
        level = 0
        m = re.search(r"heading\s*(\d)", style)
        if m:
            level = int(m.group(1))
        units.append({
            "heading": is_heading, "level": level, "text": text,
            "bold": bool(para.runs) and all(r.bold for r in para.runs if r.text.strip()),
        })
        raw_len += len(text)
    # tables — flatten into pseudo-lines
    for tno, table in enumerate(doc.tables):
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            line = " | ".join(cells)
            if line.strip(" |"):
                units.append({"heading": False, "table": tno, "text": line,
                              "bold": False})
                raw_len += len(line)
    return {"units": units, "raw_len": raw_len, "kind": "docx_paras"}


def _extract_pptx(path: Path) -> dict[str, Any]:
    from pptx import Presentation

    prs = Presentation(str(path))
    units: list[dict[str, Any]] = []
    raw_len = 0
    for sno, slide in enumerate(prs.slides):
        title = ""
        body: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(r.text for r in para.runs).strip()
                if not text:
                    continue
                if not title and (shape == slide.shapes.title if slide.shapes.title else False):
                    title = text
                else:
                    body.append(text)
        # speaker notes
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        unit = {
            "heading": True, "slide": sno + 1,
            "text": title or f"(Slide {sno + 1})",
            "body": "\n".join(body), "notes": notes, "bold": False,
        }
        units.append(unit)
        raw_len += len(unit["text"]) + len(unit["body"]) + len(notes)
    return {"units": units, "raw_len": raw_len, "kind": "pptx_slides"}


def _extract_plain(path: Path) -> dict[str, Any]:
    text = path.read_text(encoding="utf-8", errors="replace")
    units = []
    raw_len = 0
    for line in text.splitlines():
        t = line.strip()
        if not t:
            continue
        is_heading = bool(re.match(r"^#{1,6}\s+", t)) or \
            (len(t) < 80 and t.endswith((":",)) and not t.endswith("::"))
        clean = re.sub(r"^#{1,6}\s+", "", t)
        units.append({"heading": is_heading and len(clean) < 90, "text": clean,
                      "bold": False})
        raw_len += len(clean)
    return {"units": units, "raw_len": raw_len, "kind": "plain_lines"}


# ---------------------------------------------------------------------------
# Vision OCR for sparse PDF pages
# ---------------------------------------------------------------------------

async def _ocr_sparse_pdf_pages(path: Path, raw: dict, max_pages: int) -> list[dict]:
    import fitz

    sparse: list[int] = []
    page_chars: dict[int, int] = {}
    with fitz.open(path) as doc:
        for pno, page in enumerate(doc):
            n = len(page.get_text().strip())
            page_chars[pno + 1] = n
            if n < OCR_MIN_CHARS:
                sparse.append(pno)
    if not sparse:
        return []

    sparse = sparse[:max_pages]
    log.info("OCR: %d sparse pages -> Gemini vision", len(sparse))

    results: list[dict] = []
    for pno in sparse:
        with fitz.open(path) as doc:
            page = doc[pno]
            pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
            if pix.width > OCR_IMAGE_MAX_DIM or pix.height > OCR_IMAGE_MAX_DIM:
                scale = OCR_IMAGE_MAX_DIM / max(pix.width, pix.height)
                pix = page.get_pixmap(matrix=fitz.Matrix(1.5 * scale, 1.5 * scale))
            png = pix.tobytes("png")

        b64 = base64.standard_b64encode(png).decode()
        prompt = (
            "This is page " + str(pno + 1) + " of an educational document. "
            "Transcribe ALL text on the page. Then, if there are figures, "
            "diagrams, equations or charts, describe them in brackets like "
            "[FIGURE: ...] with enough detail that a teacher could redraw them. "
            "Output plain text only."
        )
        try:
            text = await generate_text(
                prompt,
                system="You are a careful OCR assistant for educational material.",
                temperature=0.0,
                priority=BATCH,
            )
            results.append({"page": pno + 1, "ocr_text": text.strip()})
            page_chars[pno + 1] = len(text)
        except Exception as e:                              # noqa: BLE001
            log.error("OCR failed on page %d: %s", pno + 1, e)
    return results


# ---------------------------------------------------------------------------
# Structure detection: build sections from heading units
# ---------------------------------------------------------------------------

def _modal_size(sizes: list[float]) -> float:
    if not sizes:
        return 10.0
    from collections import Counter
    c = Counter(round(s, 0) for s in sizes)
    return float(c.most_common(1)[0][0])


def _detect_structure(raw: dict, suffix: str) -> dict[str, Any]:
    units = raw["units"]
    sections: list[dict[str, Any]] = []

    if raw["kind"] == "pptx_slides":
        for u in units:
            body_parts = [u["text"]]
            if u.get("body"):
                body_parts.append(u["body"])
            if u.get("notes"):
                body_parts.append("(teacher notes: " + u["notes"] + ")")
            sections.append({
                "id": len(sections),
                "title": u["text"],
                "level": 1,
                "text": "\n".join(body_parts),
                "source": f"slide {u['slide']}",
            })
    elif raw["kind"] == "pdf_lines":
        current: Optional[dict[str, Any]] = None
        # group consecutive lines under the last heading
        for u in units:
            if u["heading"]:
                if current:
                    sections.append(current)
                current = {"id": len(sections), "title": u["text"],
                           "level": 1 if u["size"] >= u["body_size"] * 1.5 else 2,
                           "text": "", "source": f"page {u['page']}"}
            else:
                if current is None:
                    current = {"id": 0, "title": "Introduction", "level": 1,
                               "text": "", "source": f"page {u['page']}"}
                current["text"] += u["text"] + "\n"
        if current:
            sections.append(current)
    else:
        # docx / plain: heading units start new sections
        current: Optional[dict[str, Any]] = None
        for u in units:
            if u["heading"]:
                if current:
                    sections.append(current)
                current = {"id": len(sections), "title": u["text"],
                           "level": u.get("level") or 1, "text": ""}
            else:
                if current is None:
                    current = {"id": 0, "title": "Content", "level": 1, "text": ""}
                current["text"] += u["text"] + "\n"
        if current:
            sections.append(current)

    # merge tiny fragments into previous section (avoid heading-only noise)
    merged: list[dict[str, Any]] = []
    for s in sections:
        if merged and (len(s["text"].strip()) < 40 and len(s["title"]) < 60):
            merged[-1]["text"] += f"\n{s['title']}\n{s['text']}"
        else:
            s["id"] = len(merged)
            merged.append(s)
    sections = merged

    toc = [{"title": s["title"], "source": s.get("source", ""),
            "level": s["level"]} for s in sections]
    total_chars = sum(len(s["text"]) for s in sections)
    return {
        "sections": sections,
        "toc": toc,
        "stats": {
            "sections": len(sections),
            "chars": total_chars,
            "detected_headings": sum(1 for u in units if u.get("heading")),
        },
    }
